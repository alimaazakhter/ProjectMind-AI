"""
exports.py — Export helpers for PDF, DOCX, PPTX, and Markdown.

Each ``build_*`` function accepts a ``dict[str, str]`` of sections
and returns the file bytes (or a string for markdown).
"""

from __future__ import annotations


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------
def _sanitize(text: str) -> str:
    """Strip everything outside printable ASCII for simple PDF compatibility."""
    return text.encode("ascii", "ignore").decode("ascii").strip()


def _clean_markdown(text: str) -> str:
    """Strip common markdown formatting for plain-text exports."""
    return (
        text.replace("**", "")
        .replace("```", "")
        .replace("`", "")
        .replace("###", "")
        .replace("##", "")
        .replace("#", "")
    )


# ---------------------------------------------------------------------------
# PDF (fpdf2)
# ---------------------------------------------------------------------------
def build_pdf(sections: dict[str, str]) -> bytes:
    """Build a PDF from sections dict. Requires ``fpdf2``."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(
        0, 12,
        "AI Auto Project Generator - Blueprint",
        new_x="LMARGIN", new_y="NEXT", align="C",
    )
    pdf.ln(6)

    for title, body in sections.items():
        safe_title = _sanitize(title) or "Section"
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, safe_title, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 11)
        clean = _sanitize(_clean_markdown(body))

        for line in clean.split("\n"):
            if not line.strip():
                pdf.ln(3)
                continue
            try:
                pdf.multi_cell(0, 6, line)
            except Exception:
                pass
        pdf.ln(4)

    return bytes(pdf.output())


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------
def build_docx(sections: dict[str, str]) -> bytes:
    """Build a DOCX from sections dict. Requires ``python-docx``."""
    import io
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run("AI Auto Project Generator — Blueprint")
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)

    doc.add_paragraph()  # spacer

    for section_title, body in sections.items():
        heading = doc.add_heading(section_title, level=2)
        for run in heading.runs:
            run.font.color.rgb = RGBColor(0x6D, 0x28, 0xD9)

        clean = _clean_markdown(body)
        for line in clean.split("\n"):
            line = line.strip()
            if not line:
                continue
            para = doc.add_paragraph(line)
            para.style.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------
def build_pptx(sections: dict[str, str]) -> bytes:
    """Build a PPTX from sections dict. Requires ``python-pptx``."""
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "AI Auto Project Generator"
    slide.placeholders[1].text = "Complete Project Blueprint"

    # Content slides — one per section
    slide_layout = prs.slide_layouts[1]  # Title and Content
    for section_title, body in sections.items():
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = section_title

        clean = _clean_markdown(body)
        # Truncate to fit on a slide
        lines = [l.strip() for l in clean.split("\n") if l.strip()]
        content = "\n".join(lines[:20])  # max ~20 lines per slide
        if len(lines) > 20:
            content += "\n..."

        tf = slide.placeholders[1].text_frame
        tf.text = content
        for para in tf.paragraphs:
            para.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Markdown / README
# ---------------------------------------------------------------------------
def build_markdown(sections: dict[str, str]) -> str:
    """Build a Markdown/README string from sections dict."""
    lines = ["# AI Auto Project Generator — Blueprint\n"]
    for title, body in sections.items():
        lines.append(f"## {title}\n")
        lines.append(body.strip())
        lines.append("")  # blank line between sections
    return "\n".join(lines)
